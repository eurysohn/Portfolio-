import io
import json
import pickle
import re
import urllib.request
from pathlib import Path
from typing import Dict, List

from pdfminer.high_level import extract_pages
from pdfminer.layout import LTTextContainer
from pptx import Presentation
from docx import Document
from sklearn.feature_extraction.text import TfidfVectorizer

from data.build_synthetic_docs import generate_synthetic_docs


BASE_DIR = Path(__file__).resolve().parent
DATA_DIR = BASE_DIR / "data"
DOCS_DIR = DATA_DIR / "enterprise_knowledge"
VECTOR_DIR = BASE_DIR / "storage" / "vector_db"
INDEX_PATH = VECTOR_DIR / "index.pkl"
SEED_URLS_PATH = DATA_DIR / "seed_urls.json"
DICTIONARY_PATH = DATA_DIR / "scm_dictionary.json"
RAW_DIR = DATA_DIR / "raw_data"


def _normalize_text(text: str) -> str:
    return re.sub(r"\s+", " ", text).strip()


def _chunk_text(text: str, min_words: int = 600, max_words: int = 900) -> List[str]:
    words = text.split()
    if not words:
        return []
    chunks = []
    start = 0
    step = max_words
    while start < len(words):
        end = min(start + step, len(words))
        chunk = " ".join(words[start:end])
        chunks.append(chunk)
        start = end
    return chunks


def _load_seed_urls() -> List[str]:
    if not SEED_URLS_PATH.exists():
        return []
    try:
        data = json.loads(SEED_URLS_PATH.read_text(encoding="utf-8"))
        if isinstance(data, list):
            return [str(u) for u in data if str(u).strip()]
    except json.JSONDecodeError:
        return []
    return []


def _fetch_url_bytes(url: str) -> tuple[bytes, str]:
    try:
        request = urllib.request.Request(url, headers={"User-Agent": "scm-agent-ai-example/1.0"})
        with urllib.request.urlopen(request, timeout=15) as response:
            content_type = response.headers.get("Content-Type", "")
            raw = response.read()
    except Exception:
        return b"", ""
    return raw, content_type


def _guess_extension(url: str, content_type: str) -> str:
    url_lower = url.lower()
    if "pdf" in content_type or url_lower.endswith(".pdf"):
        return ".pdf"
    if "presentation" in content_type or url_lower.endswith(".pptx"):
        return ".pptx"
    if "word" in content_type or url_lower.endswith(".docx"):
        return ".docx"
    return ".html"


def _strip_html(text: str) -> str:
    text = re.sub(r"<script.*?>.*?</script>", " ", text, flags=re.S | re.I)
    text = re.sub(r"<style.*?>.*?</style>", " ", text, flags=re.S | re.I)
    text = re.sub(r"<[^>]+>", " ", text)
    return _normalize_text(text)


def _extract_pdf_pages(data: bytes) -> List[str]:
    pages = []
    for page_layout in extract_pages(io.BytesIO(data)):
        texts = []
        for element in page_layout:
            if isinstance(element, LTTextContainer):
                texts.append(element.get_text())
        page_text = _normalize_text(" ".join(texts))
        if page_text:
            pages.append(page_text)
    return pages


def _extract_pptx_slides(data: bytes) -> List[str]:
    slides = []
    presentation = Presentation(io.BytesIO(data))
    for slide in presentation.slides:
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
        slide_text = _normalize_text(" ".join(texts))
        if slide_text:
            slides.append(slide_text)
    return slides


def _extract_docx_text(data: bytes) -> str:
    doc = Document(io.BytesIO(data))
    texts = [para.text for para in doc.paragraphs if para.text.strip()]
    return _normalize_text(" ".join(texts))


def _write_docs(docs: List[Dict[str, str]]) -> None:
    DOCS_DIR.mkdir(parents=True, exist_ok=True)
    for doc in docs:
        filename = f"{doc['id']}.txt"
        path = DOCS_DIR / filename
        path.write_text(doc["text"], encoding="utf-8")


def _load_docs() -> List[Dict[str, str]]:
    docs = []
    for path in sorted(DOCS_DIR.glob("*.txt")):
        text = path.read_text(encoding="utf-8")
        page_text = text if ("_page_" in path.stem or "_slide_" in path.stem) else ""
        docs.append(
            {
                "id": path.stem,
                "source": str(path),
                "text": text,
                "page_text": page_text,
            }
        )
    return docs


def _build_dictionary_if_missing() -> None:
    if DICTIONARY_PATH.exists():
        return
    DICTIONARY_PATH.write_text("[]", encoding="utf-8")


def _build_vector_index(docs: List[Dict[str, str]]) -> None:
    chunks = []
    for doc in docs:
        for idx, chunk in enumerate(_chunk_text(doc["text"])):
            chunks.append(
                {
                    "chunk_id": f"{doc['id']}_chunk_{idx}",
                    "source": doc["source"],
                    "text": chunk,
                    "page_text": doc.get("page_text", ""),
                }
            )

    texts = [c["text"] for c in chunks]
    if not texts:
        raise RuntimeError("No documents available to build the vector index.")
    vectorizer = TfidfVectorizer(stop_words="english")
    matrix = vectorizer.fit_transform(texts)

    VECTOR_DIR.mkdir(parents=True, exist_ok=True)
    with INDEX_PATH.open("wb") as f:
        pickle.dump({"vectorizer": vectorizer, "matrix": matrix, "chunks": chunks}, f)


def main() -> None:
    DOCS_DIR.mkdir(parents=True, exist_ok=True)
    for old in DOCS_DIR.glob("*.txt"):
        old.unlink()

    synthetic_docs = generate_synthetic_docs()
    _write_docs(synthetic_docs)

    seed_urls = _load_seed_urls()
    if seed_urls:
        RAW_DIR.mkdir(parents=True, exist_ok=True)
        for i, url in enumerate(seed_urls, start=1):
            data, content_type = _fetch_url_bytes(url)
            if not data:
                continue
            url_lower = url.lower()
            ext = _guess_extension(url, content_type)
            raw_path = RAW_DIR / f"seed_url_{i}{ext}"
            raw_path.write_bytes(data)

            if ext == ".pdf":
                pages = _extract_pdf_pages(data)
                for page_num, page_text in enumerate(pages, start=1):
                    doc_id = f"seed_url_{i}_page_{page_num}"
                    (DOCS_DIR / f"{doc_id}.txt").write_text(page_text, encoding="utf-8")
            elif ext == ".pptx":
                slides = _extract_pptx_slides(data)
                for slide_num, slide_text in enumerate(slides, start=1):
                    doc_id = f"seed_url_{i}_slide_{slide_num}"
                    (DOCS_DIR / f"{doc_id}.txt").write_text(slide_text, encoding="utf-8")
            elif ext == ".docx":
                text = _extract_docx_text(data)
                if text:
                    doc_id = f"seed_url_{i}_page_1"
                    (DOCS_DIR / f"{doc_id}.txt").write_text(text, encoding="utf-8")
            else:
                text = _strip_html(data.decode("utf-8", errors="ignore"))
                if text:
                    doc_id = f"seed_url_{i}"
                    (DOCS_DIR / f"{doc_id}.txt").write_text(text, encoding="utf-8")

    _build_dictionary_if_missing()
    docs = _load_docs()
    _build_vector_index(docs)
    print(f"Built vector index with {len(docs)} docs at {INDEX_PATH}")


if __name__ == "__main__":
    main()
